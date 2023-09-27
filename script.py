import os
import requests
import shutil
import openpyxl
from datetime import datetime
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.util import Inches

def get_current_version():
    # Az alkalmazás beépített verziószámának lekérése
    return "1.0"

def get_latest_version():
    # A legfrissebb verziószám letöltése a megbízható forrásból (pl. egy szerverről)
    version_data = requests.get("https://raw.githubusercontent.com/akosblack/ppt_script/main/latest_version.txt").text
    return str(version_data)

def replace_text_in_ppt(pptx_file, keywords, replacements, charts):
    print("Az adatok kicserélése folyamatban...")
    # Másolja az eredeti fájlt az "output" mappába
    output_folder = "output"
    os.makedirs(output_folder, exist_ok=True)
    time = datetime.now().strftime("%Y-%m-%d")
    pptx_file_output = os.path.join(output_folder, pptx_file.replace('input\\',''))
    pptx_file_output = pptx_file_output.replace('.pptx', f'_{time}.pptx')
    # print(pptx_file_output)
    shutil.copy(pptx_file, pptx_file_output)
    # Nyissa meg a másolt fájlt
    prs = Presentation(pptx_file_output)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart:
                    chart = shape.chart
                    update_chart(chart, charts[0])
                    if len(charts) > 0:                        
                        charts.remove(charts[0])
                        print(f'charts: {charts}')
            if hasattr(shape, "text"):
                for keyword, replacement in zip(keywords, replacements):
                    shape.text = shape.text.replace(keyword, replacement)
    prs.save(pptx_file_output)
    return pptx_file_output

def update_chart(chart, chart_data):
    print(chart.chart_type)

    chart.replace_data(chart_data)

def read_input(excel_input):
    keywords = []
    replacements = []
    chart_keywords = []
    chart_datas = []
    charts = []
    # Nyissa meg az xlsx fájlt
    workbook = openpyxl.load_workbook(excel_input)
    worksheet = workbook.active
    current_chart_type = None
    # Olvassa be a tartalmat soronként
    for row in worksheet.iter_rows(values_only=True):
        if row[0] is not None:
            cell_value = str(row[0]).strip().lower()
            if any(cell_value.startswith(chart_type) for chart_type in ["column", "line", "pie", "bar", "area"]):
                print(f'cell_value in current chart type: {cell_value}')
                current_chart_type = cell_value
                chart_data = ChartData()
                if current_chart_type == 'column':
                    print(f'found chart data in excel which is column')
                    if 'title' in cell_value:
                        print(f'title: {row[1]}')
                        chart_data.chart.title.text_frame.text = row[1]
                    elif 'categories' in cell_value:
                        print(f'categories: {row[1:]}')
                        chart_data.categories = row[1:]

                    elif 'series' in cell_value:
                        print(f'series: {row[0]} {row[1]}')
                        chart_data.add_series(row[0], row[1])

                if current_chart_type is not None:
                    charts.append(chart_data)
                    #print(f'charts: {charts}')                

                chart_keyword, chart_data2 = row
                chart_keywords.append(str(chart_keyword))
                chart_datas.append(str(chart_data2))
                current_chart_type = None
                continue
            else:
                keyword, replacement = row
                keywords.append(str(keyword))
                replacements.append(str(replacement))

    return keywords, replacements, charts

if __name__ == "__main__":
    # print(get_current_version(), get_latest_version())
    if get_current_version() != get_latest_version():
        print("Frissítés elérhető. Kérlek használd a közös mappában lévő exe-t...")
        input()
        exit(1)
    try:      
        input_folder = "input"
        os.makedirs(input_folder, exist_ok=True)  
        # Keresse meg az input mappában található .xls vagy .xlsx fájlt
        excel_input = None
        for filename in os.listdir(input_folder):
            if filename.endswith((".xls", ".xlsx")):
                excel_input = os.path.join(input_folder, filename)
                break
        if excel_input is None:
            raise Exception("Nem található .xls vagy .xlsx fájl az input mappában.")

        # Keresse meg az input mappában található .ppt vagy .pptx fájlt
        template_pptx = None
        for filename in os.listdir(input_folder):
            if filename.endswith((".ppt", ".pptx")):
                template_pptx = os.path.join(input_folder, filename)
                break
        if template_pptx is None:
            raise Exception("Nem található .ppt vagy .pptx fájl az input mappában.")

        keywords, replacements, charts = read_input(excel_input)
        # print(f'chart_keywords: {chart_keywords}')
        # print(f'chart_data: {chart_datas}')
        print(keywords)
        print(replacements)
        new_pptx_file = replace_text_in_ppt(template_pptx, keywords, replacements, charts)
        os.system(f"start {new_pptx_file}")
        # Kiírja, hogy végzett
        print("Az adatok kicserelése befejezve. Nyomj meg egy billentyűt a kilépéshez...")
        input()
    except Exception as e:        
        print("Hiba történt:", e)
        print("Nyomj meg egy billentyűt a kilépéshez...")
        input()
